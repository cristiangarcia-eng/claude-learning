#!/usr/bin/env python3
"""
Create a polished single-slide PPT matching Claude10x landing page style.

Design tokens from the landing:
- Background: ~#1a1a1a (oklch 0.13)
- Brand green: #22C55E
- Card bg: ~#2b2b2b (oklch 0.18)
- Text: #fafafa
- Font headings: Outfit
- Font body: Plus Jakarta Sans
- Border radius: rounded
- Green glow: rgba(34, 197, 94, 0.12)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

ROOT = Path(__file__).parent
OUTPUT_PATH = ROOT / "screenshot-vs-text.pptx"

# ── Design tokens (Claude10x) ──
BG = (18, 18, 22)
CARD_BG = (38, 38, 44)
CARD_BG_HEX = "26262c"
GREEN = (34, 197, 94)       # #22C55E
GREEN_HEX = "22C55E"
RED = (239, 68, 68)         # #EF4444
RED_HEX = "EF4444"
TEXT_PRIMARY = (250, 250, 250)
TEXT_SECONDARY = (160, 160, 165)
TEXT_MUTED = (100, 100, 105)

FONT_HEADING = "Outfit"
FONT_BODY = "Plus Jakarta Sans"

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(*BG)


def emu(inches):
    return int(inches * 914400)


def add_text(left, top, width, height, text, size=18,
             color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.CENTER,
             font=FONT_BODY):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_card(left, top, width, height, border_hex, glow_hex):
    """Add a card with glow effect and border."""
    # Glow behind card
    glow_margin = emu(0.2)
    glow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left - glow_margin), Emu(top - glow_margin),
        Emu(width + glow_margin * 2), Emu(height + glow_margin * 2)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(*BG)
    glow.line.fill.background()
    # Outer shadow as glow
    spPr = glow._element.find(qn('p:spPr'))
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '600000')  # big blur for glow
    outerShdw.set('dist', '0')
    outerShdw.set('dir', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', glow_hex)
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', '25000')

    # Card itself
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*CARD_BG)
    card.line.color.rgb = RGBColor(int(border_hex[0:2], 16), int(border_hex[2:4], 16), int(border_hex[4:6], 16))
    card.line.width = Pt(1.5)

    # Corner radius
    spPr2 = card._element.find(qn('p:spPr'))
    prstGeom = spPr2.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        for child in list(avLst):
            avLst.remove(child)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 6000')

    return card


def add_pill(left, top, width, height, fill_rgb, text, text_size=18):
    """Add a rounded pill/badge."""
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(*fill_rgb)
    pill.line.fill.background()

    # Very round corners
    spPr = pill._element.find(qn('p:spPr'))
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        for child in list(avLst):
            avLst.remove(child)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 50000')  # max roundness

    add_text(left, top + emu(0.05), width, height,
             text, size=text_size, bold=True, color=(255, 255, 255))


# ═══════════════════════════════════════════
# LAYOUT
# ═══════════════════════════════════════════

card_w = emu(6.4)
card_h = emu(5.4)
gap = emu(0.7)
total = card_w * 2 + gap
sx = (prs.slide_width - total) // 2
card_top = emu(2.4)
lx = sx
rx = sx + card_w + gap

# ── Title ──
add_text(emu(1), emu(0.45), emu(14), emu(0.9),
         "Mismo mensaje. Dos formas de mandárselo a tu AI.",
         size=38, bold=True, font=FONT_HEADING)

# ── Subtitle ──
add_text(emu(1), emu(1.3), emu(14), emu(0.5),
         "¿Cuál usas tú?",
         size=20, color=TEXT_SECONDARY)

# ══════════════════════════
# LEFT CARD — Screenshot (red)
# ══════════════════════════
add_card(lx, card_top, card_w, card_h, RED_HEX, RED_HEX)

# Icon — simple circle with dot (camera lens)
icon_cx = lx + card_w // 2
icon_cy = card_top + emu(0.6)
# Outer ring
ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
    Emu(icon_cx - emu(0.3)), Emu(icon_cy - emu(0.3)),
    Emu(emu(0.6)), Emu(emu(0.6)))
ring.fill.background()
ring.line.color.rgb = RGBColor(*RED)
ring.line.width = Pt(2.5)
# Inner dot
dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
    Emu(icon_cx - emu(0.1)), Emu(icon_cy - emu(0.1)),
    Emu(emu(0.2)), Emu(emu(0.2)))
dot.fill.solid()
dot.fill.fore_color.rgb = RGBColor(*RED)
dot.line.fill.background()

# Label
add_text(lx, card_top + emu(1.1), card_w, emu(0.5),
         "Mandar un screenshot", size=22, bold=True, font=FONT_HEADING)

# Big number
add_text(lx, card_top + emu(1.7), card_w, emu(1.0),
         "1,585", size=72, bold=True, color=RED, font=FONT_HEADING)

# "tokens"
add_text(lx, card_top + emu(2.65), card_w, emu(0.4),
         "tokens", size=18, color=RED)

# Separator line
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(lx + emu(1.2)), Emu(card_top + emu(3.2)),
    Emu(card_w - emu(2.4)), Emu(emu(0.015)))
sep.fill.solid()
sep.fill.fore_color.rgb = RGBColor(70, 70, 75)
sep.line.fill.background()

# Bullets
bullets_left = [
    "La AI tiene que 'leer' la imagen",
    "Puede confundir caracteres",
    "No puede buscar en tus documentos",
]
y = card_top + emu(3.5)
for b in bullets_left:
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Emu(lx + emu(0.7)), Emu(y + emu(0.1)),
        Emu(emu(0.1)), Emu(emu(0.1)))
    d.fill.solid()
    d.fill.fore_color.rgb = RGBColor(*RED)
    d.line.fill.background()
    add_text(lx + emu(1.0), y, emu(4.8), emu(0.35),
             b, size=14, color=TEXT_SECONDARY, align=PP_ALIGN.LEFT)
    y += emu(0.45)


# ══════════════════════════
# RIGHT CARD — Copy-paste (green)
# ══════════════════════════
add_card(rx, card_top, card_w, card_h, GREEN_HEX, GREEN_HEX)

# Icon — three horizontal lines (text/clipboard)
for i in range(3):
    line_y = card_top + emu(0.45) + i * emu(0.2)
    line_w = emu(0.45) if i != 1 else emu(0.35)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(rx + card_w // 2 - emu(0.22)), Emu(line_y),
        Emu(line_w), Emu(emu(0.06)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*GREEN)
    bar.line.fill.background()

# Label
add_text(rx, card_top + emu(1.1), card_w, emu(0.5),
         "Copiar y pegar el texto", size=22, bold=True, font=FONT_HEADING)

# Big number
add_text(rx, card_top + emu(1.7), card_w, emu(1.0),
         "458", size=72, bold=True, color=GREEN, font=FONT_HEADING)

# "tokens"
add_text(rx, card_top + emu(2.65), card_w, emu(0.4),
         "tokens", size=18, color=GREEN)

# Separator
sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(rx + emu(1.2)), Emu(card_top + emu(3.2)),
    Emu(card_w - emu(2.4)), Emu(emu(0.015)))
sep2.fill.solid()
sep2.fill.fore_color.rgb = RGBColor(70, 70, 75)
sep2.line.fill.background()

# Bullets
bullets_right = [
    "La AI lee cada carácter perfecto",
    "Cero confusiones",
    "Puede buscar directamente en tu código",
]
y = card_top + emu(3.5)
for b in bullets_right:
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Emu(rx + emu(0.7)), Emu(y + emu(0.1)),
        Emu(emu(0.1)), Emu(emu(0.1)))
    d.fill.solid()
    d.fill.fore_color.rgb = RGBColor(*GREEN)
    d.line.fill.background()
    add_text(rx + emu(1.0), y, emu(4.8), emu(0.35),
             b, size=14, color=TEXT_SECONDARY, align=PP_ALIGN.LEFT)
    y += emu(0.45)


# ── Badge ──
pill_w = emu(4.8)
pill_h = emu(0.55)
pill_x = (prs.slide_width - pill_w) // 2
pill_y = card_top + card_h + emu(0.4)

add_pill(pill_x, pill_y, pill_w, pill_h, RED, "El screenshot cuesta 3.5x más", text_size=19)

# ── Footer ──
add_text(emu(1), emu(8.55), emu(14), emu(0.35),
         "Medido con la API de Claude  ·  Mismo prompt, mismo modelo",
         size=11, color=TEXT_MUTED)


prs.save(str(OUTPUT_PATH))
print(f"Saved: {OUTPUT_PATH}")
